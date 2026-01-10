"""
AI Chunkenizer - PowerPoint Extractor
Text extraction from .pptx files using python-pptx.
Extracts text from slides, speaker notes, and tables.

Copyright (c) 2025 VIEWVALUE LLC. All rights reserved.
Licensed under the MIT License. See LICENSE file for details.
"""

from pptx import Presentation


def extract_pptx(file_path: str) -> str:
    """
    Extract text from a PowerPoint presentation.

    Extracts text from slides, shapes, tables, and speaker notes.
    Each slide is clearly marked for context.

    Args:
        file_path: Path to the .pptx file

    Returns:
        Text from all slides with structure markers

    Raises:
        Exception: If presentation cannot be processed
    """
    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"[Slide {slide_num}]"]

            # Extract from all shapes
            for shape in slide.shapes:
                # Text content
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

                # Tables
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_rows.append(" | ".join(cells))

                    if table_rows:
                        table_text = "\n".join(table_rows)
                        slide_content.append(f"\n[Table]\n{table_text}")

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_content.append(f"\n[Notes]\n{notes}")

            text_parts.append("\n".join(slide_content))

        return "\n\n---\n\n".join(text_parts)

    except Exception as e:
        raise Exception(f"Failed to extract PowerPoint: {str(e)}")
